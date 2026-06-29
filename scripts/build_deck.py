"""Generate the Preflight presentation deck (.pptx) from the committed content.

    python scripts/build_deck.py   ->  out/Preflight_deck.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
BG = RGBColor(0x0B, 0x10, 0x20)
INK = RGBColor(0xE8, 0xED, 0xF7)
MUT = RGBColor(0x9F, 0xB0, 0xD0)
GREEN = RGBColor(0x1D, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x4D, 0x62)
ACCENT = RGBColor(0x6E, 0x9B, 0xFF)

SLIDES = [
    ("Preflight", "Prove high-stakes plans before they touch a person.",
     ["UiPath AgentHack 2026  ·  Track 3 — UiPath Test Cloud  ·  Team Preflight",
      "Branch reality. Break the plan. Prove the outcome. Then act."]),
    ("The problem", "",
     ["Enterprises automate irreversible actions — discharge, claims, payroll, shipments.",
      "They execute optimistically. A hidden dependency fails → discovered AFTER harm.",
      "Software is tested before production. Human-impacting plans are not."]),
    ("The idea", "",
     ["Test the planned real-world OUTCOME of a live case — before acting —",
      "and gate release behind a human.",
      "Not testing the software. Testing the world the plan depends on."]),
    ("How it works", "",
     ["Live case → compile policy into outcome obligations",
      "→ generate a Test Manager test case per obligation (Test Cloud)",
      "→ adversarial agent finds cross-system failure cascades",
      "→ prove + evidence → GATE: fail → human remedy → re-test ↺ ; pass → release → learn"]),
    ("Demo — caught & blocked", "",
     ["Patient looks ready. Preflight: BLOCKED, 88%.",
      "Causal chain: transport 17:15 vs pharmacy closes 17:00 → medication never collectable.",
      "A three-system cascade no checklist catches."]),
    ("Test Cloud is the heart", "",
     ["An agent turned the discharge policy into 9 executable Test Manager test cases.",
      "'Evaluate requirements and turn them into meaningful test scenarios' — automatically.",
      "The failing one is the medication-collectable test."]),
    ("Human gate + proven safe", "",
     ["Critical fail → withhold discharge → nurse approves an operational remedy (24h pharmacy).",
      "Re-test → 100%, GREEN, Released + reusable regression scenarios learned.",
      "We don't predict safety. We prove every dependency survived rehearsal."]),
    ("Built on UiPath — a blend of agents", "",
     ["Test Manager (Test Cloud) — agentic test-case generation",
      "Data Fabric — live system of record (no hardcoded data)",
      "Agent Builder — published low-code Obligation Compiler agent",
      "Coded agent via UiPath for Coding Agents (Claude Code) — the rehearsal engine  [bonus]",
      "External framework — LangGraph red-team agent, governed by UiPath",
      "UiPath Apps + live [i]-help dashboard · Identity/External Apps OAuth"]),
    ("Impact & roadmap", "",
     ["Wedge: hospital discharge (CMS readmissions, AHRQ RED).",
      "Generalizes to claims settlement, payroll, emergency payments, supply-chain release.",
      "Next: Maestro BPMN + Action Center gate; report results to Test Cloud via Orchestrator."]),
    ("Preflight", "The missing layer between a prototype on a laptop and a plan you can trust with a patient.",
     ["Repo (MIT): github.com/usv240/preflight",
      "Synthetic data · never diagnoses or prescribes · a clinician approves every change."]),
]


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _txt(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return p


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, (title, subtitle, bullets) in enumerate(SLIDES):
        s = prs.slides.add_slide(blank)
        _bg(s)
        # accent bar
        bar = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(0.12), Inches(0.7))
        bar.fill  # noqa
        # title
        tb = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.4))
        tf = tb.text_frame; tf.word_wrap = True
        big = i in (0, len(SLIDES) - 1)
        _txt(tf, title, 40 if big else 30, INK if not big else ACCENT, bold=True)
        if subtitle:
            _txt(tf, subtitle, 20 if big else 18, MUT)
        # bullets
        bb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(4.6))
        bf = bb.text_frame; bf.word_wrap = True
        for j, b in enumerate(bullets):
            color = INK
            if "BLOCKED" in b or "never collectable" in b or "fail" in b.lower():
                color = RED
            if "GREEN" in b or "Released" in b or "prove every" in b.lower():
                color = GREEN
            p = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
            p.space_after = Pt(10)
            r = p.add_run(); r.text = ("•  " if len(bullets) > 1 and not big else "") + b
            r.font.size = Pt(20 if not big else 18); r.font.color.rgb = color; r.font.name = "Segoe UI"
        # footer
        ft = s.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
        _txt(ft.text_frame, "Preflight · UiPath AgentHack 2026 · Track 3 — Test Cloud", 10, MUT)

    out = ROOT / "out"; out.mkdir(exist_ok=True)
    path = out / "Preflight_deck.pptx"
    prs.save(str(path))
    print(f"Wrote {path}  ({len(SLIDES)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
